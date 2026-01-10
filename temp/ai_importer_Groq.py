import os
import time
import json
import shutil
import psycopg2
from psycopg2.extras import Json
import pandas as pd
import docx
from pypdf import PdfReader # Thư viện đọc PDF
from groq import Groq # Thư viện AI mới
from datetime import datetime, date
from pptx import Presentation # Thư viện đọc PowerPoint
import win32com.client # Thư viện đọc .mpp qua COM

# --- CẤU HÌNH HỆ THỐNG ---

# 1. Cấu hình Groq API
GROQ_API_KEY = "" # <--- DÁN KEY GROQ CỦA BẠN VÀO ĐÂY
client = Groq(api_key=GROQ_API_KEY)

# Sử dụng Model Llama 3.3 70B mới nhất
MODEL_NAME = "llama-3.3-70b-versatile" 

# 2. Cấu hình PostgreSQL (IP NAS của bạn)
DB_CONFIG = {
    "dbname": "PMIS",
    "user": "postgres",
    "password": "O*&-Unh-LNG-%^#",
    "host": "100.94.213.83",
    "port": "2345"
}

# Cấu hình Đường dẫn
SOURCE_FOLDER = r"C:\Temp"
PROCESSED_FOLDER = r"C:\Temp\01_Processed_AI"
ERROR_FOLDER = r"C:\Temp\99_Error"

# ĐÃ XÓA BỎ DÒNG ALLOWED_EXTENSIONS ĐỂ KHÔNG GIỚI HẠN FILE

# Tạo folder nếu chưa có
os.makedirs(PROCESSED_FOLDER, exist_ok=True)
os.makedirs(ERROR_FOLDER, exist_ok=True)

# --- CÁC HÀM CƠ SỞ DỮ LIỆU (GIỮ NGUYÊN) ---

def get_db_connection():
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        return conn
    except Exception as e:
        print(f"Lỗi kết nối Database: {e}")
        return None

def create_table_if_not_exists():
    conn = get_db_connection()
    if not conn: return
    commands = ("""CREATE TABLE IF NOT EXISTS tbl_documents (id SERIAL PRIMARY KEY, file_name VARCHAR(255), doc_type VARCHAR(50), doc_date DATE, project_name VARCHAR(255), sender VARCHAR(255), receiver VARCHAR(255), summary TEXT, keywords TEXT[], raw_content TEXT, ai_analysis_json JSONB, processed_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP, file_path_original TEXT)""",)
    cur = conn.cursor()
    try:
        for command in commands: cur.execute(command)
        conn.commit()
    except (Exception, psycopg2.DatabaseError) as error: print(f"Lỗi tạo bảng: {error}")
    finally:
        if conn: conn.close()

def save_to_postgres(file_name, file_path, data):
    conn = get_db_connection()
    if not conn: return False
    sql = """INSERT INTO tbl_documents (file_name, doc_type, doc_date, project_name, sender, receiver, summary, keywords, ai_analysis_json, file_path_original) VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s)"""
    try:
        cur = conn.cursor()
        # Xử lý chuẩn hóa ngày tháng
        normalized_date = normalize_doc_date(data.get('doc_date'))
        # Xử lý keywords đảm bảo là list
        keywords = data.get('keywords') if isinstance(data.get('keywords'), list) else []
        
        cur.execute(sql, (
            file_name, 
            data.get('doc_type'), 
            normalized_date, 
            data.get('project_name'), 
            data.get('sender'), 
            data.get('receiver'), 
            data.get('summary'), 
            keywords, 
            Json(data), 
            file_path
        ))
        conn.commit()
        cur.close()
        return True
    except (Exception, psycopg2.DatabaseError) as error:
        print(f"Lỗi INSERT SQL: {error}")
        conn.rollback()
        return False
    finally:
        if conn: conn.close()

def normalize_doc_date(value):
    """Chuyển ngày về kiểu date; trả None nếu không hợp lệ."""
    if value is None: return None
    if isinstance(value, date): return value
    if isinstance(value, datetime): return value.date()
    if isinstance(value, str):
        v = value.strip()
        if not v: return None
        for fmt in ("%Y-%m-%d", "%d-%m-%Y", "%Y/%m/%d", "%d/%m/%Y"):
            try: return datetime.strptime(v, fmt).date()
            except ValueError: pass
        for fmt in ("%Y-%m", "%Y/%m"):
            try: dt = datetime.strptime(v, fmt); return date(dt.year, dt.month, 1)
            except ValueError: pass
        for fmt in ("%Y",):
            try: dt = datetime.strptime(v, fmt); return date(dt.year, 1, 1)
            except ValueError: pass
    return None

# --- HÀM TRÍCH XUẤT NỘI DUNG (ĐÃ MỞ RỘNG CHO MỌI FILE) ---

def extract_text_from_file(file_path):
    """
    Thử đọc mọi loại file. 
    Nếu đọc được -> Trả về text.
    Nếu không hỗ trợ (Ảnh, Zip...) -> Trả về None.
    """
    ext = os.path.splitext(file_path)[1].lower()
    text = ""
    
    try:
        # Nhóm file PDF
        if ext == '.pdf':
            reader = PdfReader(file_path)
            for page in reader.pages:
                text += (page.extract_text() or "") + "\n"
        
        # Nhóm file Word
        elif ext in ['.docx', '.doc']:
            doc = docx.Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
            for table in doc.tables:
                for row in table.rows:
                    text += "\n" + " | ".join([cell.text for cell in row.cells])
        
        # Nhóm file Excel & CSV
        elif ext in ['.xlsx', '.xls', '.csv']:
            if ext == '.csv':
                df = pd.read_csv(file_path)
            else:
                df = pd.read_excel(file_path)
            text = df.to_markdown(index=False)
            
        # Nhóm file PowerPoint
        elif ext in ['.pptx', '.ppt']:
            try:
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            except Exception as e:
                print(f"-> Lỗi đọc PowerPoint: {e}")
                return None
        
        # Nhóm file Microsoft Project (.mpp)
        elif ext == '.mpp':
            try:
                msp = win32com.client.Dispatch("MSProject.Application")
                msp.FileOpen(file_path)
                proj = msp.ActiveProject
                
                # Trích xuất thông tin dự án
                text += f"Tên dự án: {proj.Name}\n"
                text += f"Ngày bắt đầu: {proj.ProjectStart}\n"
                text += f"Ngày kết thúc: {proj.ProjectFinish}\n\n"
                text += "=== Danh sách công việc ===\n"
                
                for task in proj.Tasks:
                    if task is not None:
                        text += f"- {task.Name} | Bắt đầu: {task.Start} | Kết thúc: {task.Finish} | Tiến độ: {task.PercentComplete}%\n"
                
                msp.FileClose()
                msp.Quit()
            except Exception as e:
                print(f"-> Lỗi đọc MS Project (cần cài MS Project): {e}")
                return None
        
        # Nhóm file Văn bản thuần / Code (Mở rộng thêm)
        elif ext in ['.txt', '.md', '.py', '.json', '.xml', '.html', '.sql', '.log']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
        
        # Nhóm không hỗ trợ trích xuất text (Ảnh, Zip, Exe...)
        else:
            print(f"-> Định dạng {ext} chưa hỗ trợ trích xuất text tự động. Bỏ qua.")
            return None

        return text

    except Exception as e:
        print(f"-> Lỗi khi cố đọc file {ext}: {e}")
        return None

# --- HÀM GỌI GROQ AI ---

def analyze_with_groq(file_path):
    print(f"Đang đọc nội dung file: {os.path.basename(file_path)}...")
    
    # 1. Bước quan trọng: Python đọc file thành chữ
    raw_text = extract_text_from_file(file_path)
    
    if not raw_text or len(raw_text.strip()) == 0:
        print("File rỗng hoặc không đọc được nội dung.")
        return None

    # Giới hạn token (Llama 3.3 70B limit 12800 tokens ~ 38400 ký tự, để an toàn cắt ở 20000)
    if len(raw_text) > 20000:
        raw_text = raw_text[:20000] + "\n...[Đã cắt bớt nội dung thừa]..."

    # 2. Tạo Prompt
    system_prompt = "Bạn là chuyên gia quản lý hồ sơ xây dựng. Nhiệm vụ của bạn là trích xuất thông tin từ văn bản người dùng cung cấp và trả về định dạng JSON."
    
    user_prompt = f"""
    Hãy phân tích nội dung văn bản dưới đây và trả về JSON (chỉ JSON, không lời dẫn).
    
    Các trường cần lấy:
    - "doc_type": Loại văn bản (Quyết định, Tờ trình, Hợp đồng, Bản vẽ, v.v.)
    - "doc_date": Ngày văn bản (Format chuẩn YYYY-MM-DD), nếu không có để null.
    - "project_name": Tên dự án.
    - "sender": Nơi gửi / Cơ quan ban hành.
    - "receiver": Nơi nhận.
    - "summary": Tóm tắt 50 từ tiếng Việt.
    - "keywords": ["từ khóa 1", "từ khóa 2"].

    NỘI DUNG VĂN BẢN:
    {raw_text}
    """

    try:
        # 3. Gửi lên Groq
        chat_completion = client.chat.completions.create(
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            model=MODEL_NAME,
            temperature=0.1, # Nhiệt độ thấp để lấy dữ liệu chính xác
            response_format={"type": "json_object"} # Bắt buộc trả về JSON
        )

        result_json = chat_completion.choices[0].message.content
        return json.loads(result_json)

    except Exception as e:
        print(f"Lỗi API Groq: {e}")
        return None

# --- MAIN (ĐÃ CẬP NHẬT LOGIC KHÔNG LỌC EXTENSION) ---

def main():
    print("--- BẮT ĐẦU (CHẾ ĐỘ KHÔNG GIỚI HẠN ĐỊNH DẠNG) ---")
    create_table_if_not_exists()
    
    # Lấy danh sách file (Bỏ qua filter extension)
    files = [f for f in os.listdir(SOURCE_FOLDER) if os.path.isfile(os.path.join(SOURCE_FOLDER, f))]
    print(f"Tìm thấy {len(files)} file.")

    for file_name in files:
        # Chỉ bỏ qua file rác hệ thống và file tạm Word
        if file_name.startswith("~$") or file_name == "Thumbs.db": 
            continue
            
        file_path = os.path.join(SOURCE_FOLDER, file_name)

        print(f"\n>>> Xử lý: {file_name}")
        
        # Gọi hàm phân tích
        ai_data = analyze_with_groq(file_path)
        
        if ai_data:
            print("Đã có dữ liệu JSON. Lưu DB...")
            saved = save_to_postgres(file_name, file_path, ai_data)
            if saved:
                shutil.move(file_path, os.path.join(PROCESSED_FOLDER, file_name))
            else:
                shutil.move(file_path, os.path.join(ERROR_FOLDER, file_name))
        else:
            print("-> Không trích xuất được thông tin (File lạ hoặc lỗi AI). Di chuyển sang Error.")
            shutil.move(file_path, os.path.join(ERROR_FOLDER, file_name))
            
        # Groq nhanh nhưng vẫn nên nghỉ 2-3s để tránh spam
        time.sleep(3)

if __name__ == "__main__":
    main()