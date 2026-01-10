import os
import json
import shutil
import time

# --- CÁC THƯ VIỆN ĐỌC FILE ---
try:
    from pypdf import PdfReader
    from docx import Document
    from openpyxl import load_workbook
    from pptx import Presentation
except ImportError as e:
    print("❌ Thiếu thư viện! Bạn hãy chạy lệnh sau để cài đặt:")
    print("pip install pypdf python-docx openpyxl python-pptx")
    exit()

# ================= CẤU HÌNH =================
# Đường dẫn chứa các file lộn xộn cần sắp xếp
SOURCE_FOLDER = r"C:\Temp"

# Đường dẫn file cấu hình (được tạo từ Bước 1)
CONFIG_FILE = "folder_map.json"
# ============================================

def load_config(filename):
    with open(filename, 'r', encoding='utf-8') as f:
        return json.load(f)

# --- CÁC HÀM ĐỌC NỘI DUNG FILE ---
def read_txt(path):
    try:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except: return ""

def read_pdf(path):
    text = ""
    try:
        reader = PdfReader(path)
        # Chỉ đọc tối đa 5 trang đầu để tiết kiệm thời gian
        for i, page in enumerate(reader.pages):
            if i > 5: break 
            text += page.extract_text() + " "
    except: pass
    return text

def read_docx(path):
    text = ""
    try:
        doc = Document(path)
        for para in doc.paragraphs:
            text += para.text + " "
    except: pass
    return text

def read_excel(path):
    text = ""
    try:
        # data_only=True để chỉ đọc giá trị (bỏ qua công thức)
        wb = load_workbook(path, data_only=True)
        
        # SỬA ĐỔI: Duyệt qua TẤT CẢ các sheet trong workbook
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            
            # Đọc nội dung từng ô trong sheet
            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    # Kiểm tra nếu cell có dữ liệu (không phải None)
                    if cell: 
                        # Chuyển cell thành string và thêm khoảng trắng
                        text += str(cell) + " "
                        
    except Exception as e:
        # print(f"Lỗi đọc Excel {path}: {e}") # Có thể bỏ comment để debug
        pass
        
    return text

def read_pptx(path):
    text = ""
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
    except: pass
    return text

def get_file_content(filepath):
    """Điều hướng file đến hàm đọc tương ứng dựa trên đuôi file"""
    ext = os.path.splitext(filepath)[1].lower()
    
    if ext == '.txt':
        return read_txt(filepath)
    elif ext == '.pdf':
        return read_pdf(filepath)
    elif ext in ['.docx', '.doc']: # .doc cũ có thể lỗi, tốt nhất là .docx
        return read_docx(filepath)
    elif ext in ['.xlsx', '.xls']:
        return read_excel(filepath)
    elif ext in ['.pptx', '.ppt']:
        return read_pptx(filepath)
    return ""

# --- HÀM CHÍNH: PHÂN LOẠI VÀ DI CHUYỂN ---
def organize_files_smart():
    if not os.path.exists(CONFIG_FILE):
        print(f"❌ Không tìm thấy file {CONFIG_FILE}. Hãy chạy code Bước 1 trước.")
        return

    folder_mapping = load_config(CONFIG_FILE)
    files_moved = 0
    
    # DEBUG: Hiển thị cấu hình folder_mapping
    print("=" * 50)
    print("📋 [DEBUG] Nội dung folder_map.json:")
    for folder, keywords in folder_mapping.items():
        print(f"   📁 {folder}")
        print(f"      🔑 Keywords: {keywords}")
    print("=" * 50)
    
    print(f"🚀 Bắt đầu quét thư mục: {SOURCE_FOLDER}")
    print("-" * 50)
    
    # DEBUG: Liệt kê tất cả file trong thư mục nguồn
    all_items = os.listdir(SOURCE_FOLDER)
    files_only = [f for f in all_items if os.path.isfile(os.path.join(SOURCE_FOLDER, f))]
    print(f"📂 [DEBUG] Tổng số item trong {SOURCE_FOLDER}: {len(all_items)}")
    print(f"📄 [DEBUG] Số file (không tính folder): {len(files_only)}")
    if files_only:
        print("   Danh sách file:")
        for f in files_only[:20]:  # Hiển thị tối đa 20 file
            print(f"      - {f}")
        if len(files_only) > 20:
            print(f"      ... và {len(files_only) - 20} file khác")
    else:
        print("   ⚠️ KHÔNG CÓ FILE NÀO trong thư mục này!")
    print("-" * 50)

    for filename in os.listdir(SOURCE_FOLDER):
        source_file_path = os.path.join(SOURCE_FOLDER, filename)
        
        if not os.path.isfile(source_file_path):
            continue
            
        file_name_lower = filename.lower()
        destination_found = None
        
        # --- CHIẾN THUẬT 1: TÌM THEO TÊN FILE (Nhanh) ---
        print(f"🔍 Đang kiểm tra: {filename}...", end="")
        
        for folder, keywords in folder_mapping.items():
            if any(k.lower() in file_name_lower for k in keywords):
                destination_found = folder
                print(" -> ✅ Khớp TÊN file")
                break
        
        # --- CHIẾN THUẬT 2: TÌM THEO NỘI DUNG (Nếu tên không khớp) ---
        if not destination_found:
            content = get_file_content(source_file_path).lower()
            if content:
                for folder, keywords in folder_mapping.items():
                    # Chỉ tìm nếu từ khóa xuất hiện trong nội dung
                    if any(k.lower() in content for k in keywords):
                        destination_found = folder
                        print(" -> ✅ Khớp NỘI DUNG file")
                        break
            
        # --- DI CHUYỂN FILE ---
        if destination_found:
            try:
                # Xử lý trùng tên
                final_path = os.path.join(destination_found, filename)
                if os.path.exists(final_path):
                    name, ext = os.path.splitext(filename)
                    timestamp = int(time.time())
                    final_path = os.path.join(destination_found, f"{name}_{timestamp}{ext}")

                shutil.move(source_file_path, final_path)
                # print(f"   🚚 Di chuyển đến: {destination_found}")
                files_moved += 1
            except Exception as e:
                print(f"\n   ❌ Lỗi di chuyển: {e}")
        else:
            print(" -> ⚠️ Không tìm thấy nơi phù hợp")

    print("-" * 50)
    print(f"🎉 Hoàn tất! Đã xử lý và di chuyển {files_moved} file.")

# --- CHẠY CHƯƠNG TRÌNH ---
if __name__ == "__main__":
    if os.path.exists(SOURCE_FOLDER):
        organize_files_smart()
    else:
        print(f"❌ Thư mục nguồn không tồn tại: {SOURCE_FOLDER}")