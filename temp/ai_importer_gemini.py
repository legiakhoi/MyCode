import os
import time
import json
import shutil
import psycopg2
from psycopg2.extras import Json
import pandas as pd
import docx
from pptx import Presentation
from pypdf import PdfReader
import google.generativeai as genai  # Thư viện Google Gemini
from datetime import datetime
import unicodedata
import tempfile

# --- CẤU HÌNH HỆ THỐNG ---

# 1. Cấu hình Google Gemini API
# Lưu ý: Đây là Key Google Gemini của bạn (đã thay thế cho Groq)
GOOGLE_API_KEY = "AIzaSyDiXw9ZierqpR3rxFU1aHgxeuLZFVrAkco" # <--- BẠN HÃY DÁN LẠI KEY GOOGLE CỦA BẠN VÀO ĐÂY (Key bắt đầu bằng AIzaSy...)
genai.configure(api_key=GOOGLE_API_KEY)

# Sử dụng Model Flash (Miễn phí 15 RPM, nhanh và hỗ trợ tiếng Việt tốt)
MODEL_NAME = "gemini-flash-latest"

# 2. Cấu hình PostgreSQL
DB_CONFIG = {
    "dbname": "PMIS",           # Đã đổi tên DB chuẩn
    "user": "postgres",
    "password": "O*&-Unh-LNG-%^#",
    "host": "100.94.213.83",    # IP NAS của bạn
    "port": "2345"
}

# 3. Cấu hình Đường dẫn (Windows - C:\Temp)
SOURCE_FOLDER = r"C:\Temp"
PROCESSED_FOLDER = r"C:\Temp\01_Processed_AI"
ERROR_FOLDER = r"C:\Temp\99_Error"

# Tạo folder nếu chưa có
os.makedirs(PROCESSED_FOLDER, exist_ok=True)
os.makedirs(ERROR_FOLDER, exist_ok=True)

# --- CÁC HÀM CƠ SỞ DỮ LIỆU ---

def remove_accents(input_str):
    """Loại bỏ dấu tiếng Việt trong chuỗi"""
    nfkd_form = unicodedata.normalize('NFKD', input_str)
    return "".join([c for c in nfkd_form if not unicodedata.combining(c)])

def get_db_connection():
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        return conn
    except Exception as e:
        print(f"Lỗi kết nối Database: {e}")
        return None

def create_table_if_not_exists():
    conn = get_db_connection()
    if not conn: return
    
    # Tạo bảng tbl_documents nếu chưa có
    commands = (
        """
        CREATE TABLE IF NOT EXISTS tbl_documents (
            id SERIAL PRIMARY KEY,
            file_name VARCHAR(255),
            doc_type VARCHAR(50),
            doc_date DATE,
            project_name VARCHAR(255),
            sender VARCHAR(255),
            receiver VARCHAR(255),
            summary TEXT,
            keywords TEXT[],
            raw_content TEXT,
            ai_analysis_json JSONB,
            processed_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            file_path_original TEXT
        )
        """,
    )
    cur = conn.cursor()
    try:
        for command in commands:
            cur.execute(command)
        conn.commit()
    except (Exception, psycopg2.DatabaseError) as error:
        print(f"Lỗi tạo bảng: {error}")
    finally:
        if conn: conn.close()

def save_to_postgres(file_name, file_path, data):
    conn = get_db_connection()
    if not conn: return False
    
    sql = """
    INSERT INTO tbl_documents 
    (file_name, doc_type, doc_date, project_name, sender, receiver, summary, keywords, ai_analysis_json, file_path_original)
    VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
    """
    try:
        cur = conn.cursor()
        # Cắt ngắn dữ liệu nếu quá dài để tránh lỗi DB
        doc_type = data.get('doc_type')
        if doc_type and len(doc_type) > 50:
            doc_type = doc_type[:50]

        cur.execute(sql, (
            file_name,
            doc_type,
            data.get('doc_date'),
            data.get('project_name'),
            data.get('sender'),
            data.get('receiver'),
            data.get('summary'),
            data.get('keywords'),
            Json(data),
            file_path
        ))
        conn.commit()
        cur.close()
        return True
    except (Exception, psycopg2.DatabaseError) as error:
        print(f"Lỗi INSERT SQL: {error}")
        conn.rollback()
        return False
    finally:
        if conn: conn.close()

# --- HÀM TRÍCH XUẤT NỘI DUNG ---

def extract_text_from_file(file_path):
    """Đọc mọi loại file thành chữ để gửi cho AI"""
    ext = os.path.splitext(file_path)[1].lower()
    text = ""
    try:
        if ext == '.pdf':
            # PDF sẽ được xử lý riêng bằng cách upload trực tiếp lên Gemini
            return "PDF_FILE" 
        elif ext == '.docx':
            doc = docx.Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
            for table in doc.tables:
                for row in table.rows:
                    text += "\n" + " | ".join([cell.text for cell in row.cells])
        elif ext == '.pptx':
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif ext in ['.xlsx', '.xls']:
            df = pd.read_excel(file_path)
            text = df.to_markdown(index=False)
        elif ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
        return text
    except Exception as e:
        print(f"Lỗi đọc file {ext}: {e}")
        return None

# --- HÀM GỌI GOOGLE GEMINI (ĐÃ SỬA LẠI HOÀN CHỈNH) ---

def analyze_with_gemini(file_path):
    print(f"Đang đọc nội dung file: {os.path.basename(file_path)}...")
    
    # 1. Đọc nội dung file thành chữ
    raw_text = extract_text_from_file(file_path)
    
    if not raw_text:
        print("File rỗng hoặc không đọc được nội dung.")
        return None

    # 2. Tạo Prompt (Câu lệnh cho AI)
    base_prompt = """
    Bạn là chuyên gia quản lý hồ sơ xây dựng. Nhiệm vụ của bạn là trích xuất thông tin từ văn bản sau và trả về kết quả dưới dạng JSON thuần túy (no markdown).
    
    Các trường cần lấy:
    - "doc_type": Loại văn bản (Quyết định, Tờ trình, Hợp đồng, Bản vẽ, v.v.)
    - "doc_date": Ngày văn bản (Format chuẩn YYYY-MM-DD), nếu không có để null.
    - "project_name": Tên dự án.
    - "sender": Nơi gửi / Cơ quan ban hành.
    - "receiver": Nơi nhận.
    - "summary": Tóm tắt nội dung (khoảng 50 từ tiếng Việt).
    - "keywords": ["từ khóa 1", "từ khóa 2"].
    """

    try:
        # 3. Gửi lên Google Gemini
        model = genai.GenerativeModel(MODEL_NAME)
        generation_config = genai.GenerationConfig(response_mime_type="application/json")
        
        # Thêm cơ chế Retry (Thử lại) khi gặp lỗi 429 (Quota Exceeded)
        max_retries = 3
        for attempt in range(max_retries):
            try:
                response = None
                
                # Xử lý riêng cho PDF (Upload file)
                if raw_text == "PDF_FILE":
                    # Tạo bản sao tạm thời để xử lý ổn định hơn (Fix lỗi mạng/tên file)
                    fd, temp_path = tempfile.mkstemp(suffix=".pdf")
                    os.close(fd)
                    
                    shutil.copy2(file_path, temp_path)

                    try:
                        # Upload từ temp_path thay vì file_path (ổ Y)
                        uploaded_file = genai.upload_file(temp_path, mime_type="application/pdf")
                        
                        # Chờ file xử lý xong
                        while uploaded_file.state.name == "PROCESSING":
                            print(".", end="", flush=True)
                            time.sleep(2)
                            uploaded_file = genai.get_file(uploaded_file.name)
                            
                        if uploaded_file.state.name == "FAILED":
                            raise ValueError("Google không đọc được file PDF này.")
                            
                        # QUAN TRỌNG: Prompt phải nằm trong list cùng với file
                        response = model.generate_content([base_prompt, uploaded_file], generation_config=generation_config)
                        uploaded_file.delete() # Xóa file trên cloud sau khi xong
                    finally:
                        # Dọn dẹp file rác
                        if os.path.exists(temp_path):
                            os.remove(temp_path)
                    
                else:
                    # Xử lý cho các file Office/Text (Gửi text)
                    full_prompt = f"{base_prompt}\n\nNỘI DUNG VĂN BẢN CẦN XỬ LÝ:\n{raw_text[:30000]}"
                    response = model.generate_content(full_prompt, generation_config=generation_config)
                
                return json.loads(response.text)
                
            except Exception as e:
                if "429" in str(e):
                    wait_time = 60 # Đợi 60 giây nếu bị Google chặn
                    print(f"Gặp lỗi Quota (429). Đang đợi {wait_time}s để thử lại (Lần {attempt+1}/{max_retries})...")
                    time.sleep(wait_time)
                else:
                    raise e # Nếu lỗi khác thì ném ra ngoài để xử lý
                    
        return None

    except Exception as e:
        print(f"Lỗi API Gemini: {e}")
        return None

# --- MAIN ---

def main():
    print("--- BẮT ĐẦU (SỬ DỤNG GOOGLE GEMINI) ---")
    create_table_if_not_exists()
    
    files = [f for f in os.listdir(SOURCE_FOLDER) if os.path.isfile(os.path.join(SOURCE_FOLDER, f))]
    print(f"Tìm thấy {len(files)} file.")

    for file_name in files:
        if file_name.startswith("~$"): continue
        
        # Xử lý đổi tên file tiếng Việt để tránh lỗi
        clean_name = remove_accents(file_name)
        original_path = os.path.join(SOURCE_FOLDER, file_name)
        new_path = os.path.join(SOURCE_FOLDER, clean_name)
        
        if file_name != clean_name:
            try:
                os.rename(original_path, new_path)
                print(f"Đã đổi tên file: {file_name} -> {clean_name}")
                file_name = clean_name
                file_path = new_path
            except Exception as e:
                print(f"Lỗi đổi tên file {file_name}: {e}")
                file_path = original_path # Fallback
        else:
            file_path = original_path
        
        print(f"\n>>> Xử lý: {file_name}")
        
        # GỌI HÀM GEMINI
        ai_data = analyze_with_gemini(file_path)
        
        if ai_data:
            print("Đã có dữ liệu JSON. Lưu DB...")
            saved = save_to_postgres(file_name, file_path, ai_data)
            
            if saved:
                print("-> Thành công! Di chuyển sang Processed.")
                shutil.move(file_path, os.path.join(PROCESSED_FOLDER, file_name))
            else:
                print("-> Lỗi Database! Di chuyển sang Error.")
                shutil.move(file_path, os.path.join(ERROR_FOLDER, file_name))
        else:
            print("-> Lỗi phân tích AI! Di chuyển sang Error.")
            shutil.move(file_path, os.path.join(ERROR_FOLDER, file_name))
            
        # Nghỉ 15 giây giữa các file để Google không chặn (Rate Limit)
        print("Nghỉ 15 giây...")
        time.sleep(15)

if __name__ == "__main__":
    main()